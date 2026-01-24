from __future__ import annotations

import os
import re
from dataclasses import dataclass
from typing import List, Tuple
import csv
import json

from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from .config import settings


# Prefer keeping paragraph boundaries; avoid collapsing all newlines into spaces
PARA_SPLIT_RE = re.compile(r"(?:\r?\n){2,}")
SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")
UPPER_HEADING_RE = re.compile(r"^[A-Z0-9][A-Z0-9 \-:]{2,}$")
NUMBERED_HEADING_RE = re.compile(r"^(?:[IVXLCDM]+\.|\d+(?:\.\d+)*\.|[A-Z]\.)\s+.+")
PAGE_FOOTER_RE = re.compile(r"^\s*page\s+\d+(?:\s+of\s+\d+)?\s*$", re.I)


@dataclass
class ChunkParams:
    chunk_size: int = 2500
    chunk_overlap: int = 250
    # Optional custom separator order for recursive splitting
    separators: tuple[str, ...] = ("\n\n", "\n", ". ", " ", "")


def _normalize_whitespace_preserve_paragraphs(text: str) -> str:
    """Normalize whitespace but preserve blank lines as paragraph boundaries."""
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse more than two consecutive newlines to exactly two
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Normalize spaces within lines
    lines = []
    for ln in text.split("\n"):
        ln = re.sub(r"\s+", " ", ln).strip()
        lines.append(ln)
    text = "\n".join(lines)
    # Restore paragraph boundaries
    text = re.sub(r"(\n\s*){3,}", "\n\n", text)
    return text.strip()


def _fix_hyphenation(text: str) -> str:
    """Fix common PDF hyphenation like 'exam-\nple' -> 'example'."""
    # Join words broken by hyphen at line end
    text = re.sub(r"-\n(?=\w)", "", text)
    # Remove lone hyphens surrounded by newlines
    text = re.sub(r"\n-\n", "\n", text)
    # Replace single newlines inside paragraphs with spaces (but keep double newlines)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"(?<!\n)\n(?!\n)", " ", text)
    return text


def _insert_heading_boundaries(text: str) -> str:
    """Insert extra blank lines around detected headings to improve chunk boundaries."""
    out_lines: List[str] = []
    for ln in text.split("\n"):
        if UPPER_HEADING_RE.match(ln) or NUMBERED_HEADING_RE.match(ln):
            if out_lines and out_lines[-1] != "":
                out_lines.append("")
            out_lines.append(ln)
            out_lines.append("")
        else:
            out_lines.append(ln)
    return "\n".join(out_lines)


def _remove_common_headers_footers(pages: List[str]) -> List[str]:
    """Heuristic removal of repeating headers/footers across pages."""
    if not pages or len(pages) < 3:
        return pages
    # Collect first and last non-empty lines per page
    first_lines: List[str] = []
    last_lines: List[str] = []
    for p in pages:
        ls = [l.strip() for l in p.split("\n") if l.strip()]
        if not ls:
            first_lines.append("")
            last_lines.append("")
            continue
        first_lines.append(ls[0])
        last_lines.append(ls[-1])
    def most_common(cand: List[str]) -> str:
        from collections import Counter
        c = Counter([x for x in cand if x])
        return c.most_common(1)[0][0] if c else ""
    first_common = most_common(first_lines)
    last_common = most_common(last_lines)
    cleaned_pages: List[str] = []
    for p in pages:
        ls = p.split("\n")
        if first_common:
            ls = ls[1:] if ls and ls[0].strip() == first_common else ls
        if last_common:
            if ls and ls[-1].strip() == last_common:
                ls = ls[:-1]
        # Remove generic page footers like "Page X of Y"
        ls = [l for l in ls if not PAGE_FOOTER_RE.match(l.strip())]
        cleaned_pages.append("\n".join(ls))
    return cleaned_pages


def read_text_from_file(path: str) -> Tuple[str, str]:
    """
    Return (text, source_type) from a supported file.
    source_type: pdf|html|txt|docx|pptx|xlsx|xml|csv|md|json|image|audio|video
    """
    ext = os.path.splitext(path)[1].lower()
    # Documents
    if ext == ".pdf":
        return extract_text_from_pdf(path), "pdf"
    if ext in {".html", ".htm"}:
        return extract_text_from_html(path), "html"
    if ext == ".docx":
        return extract_text_from_docx(path), "docx"
    if ext == ".pptx":
        return extract_text_from_pptx(path), "pptx"
    if ext in {".xlsx", ".xls"}:
        return extract_text_from_xlsx(path), "xlsx"
    if ext in {".txt", ""}:
        return extract_text_from_txt(path), "txt"
    if ext == ".xml":
        return extract_text_from_xml(path), "xml"
    if ext == ".csv":
        return extract_text_from_csv(path), "csv"
    if ext == ".md":
        return extract_text_from_md(path), "md"
    if ext == ".json":
        return extract_text_from_json(path), "json"
    # Images (OCR)
    if ext in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif"}:
        return extract_text_from_image(path), "image"
    # Audio/Video (transcription)
    if ext in {".mp3", ".wav", ".m4a", ".flac", ".ogg"}:
        return extract_text_from_av(path, kind="audio"), "audio"
    if ext in {".mp4", ".mov", ".mkv", ".webm", ".avi"}:
        return extract_text_from_av(path, kind="video"), "video"
    # Fallback: read as text if possible
    try:
        return extract_text_from_txt(path), ext.lstrip('.') or 'txt'
    except Exception:
        raise ValueError(f"Unsupported file type: {ext}")


def extract_text_from_pdf(path: str) -> str:
    """Robust PDF extraction.
    Order of preference:
      1) PyMuPDF (if enabled): page.get_text("text"); remove common headers/footers; hyphenation fix; preserve paragraphs
      2) pypdf: page.extract_text(); hyphenation fix; preserve paragraphs
      3) pdfplumber fallback for table/figure-heavy PDFs when pypdf output is sparse
    """
    # Optional: use PyMuPDF if enabled and available for better extraction
    if getattr(settings, "use_pymupdf", False):
        try:
            import fitz  # PyMuPDF
            pages_raw: List[str] = []
            with fitz.open(path) as doc:
                for page in doc:
                    # Use textual extraction; "text" preserves reading order better than "blocks" in many docs
                    t = page.get_text("text") or ""
                    pages_raw.append(t)
            # Remove common headers/footers
            pages_clean = _remove_common_headers_footers(pages_raw)
            text = "\n\n".join(pages_clean)
            text = _fix_hyphenation(text)
            text = _normalize_whitespace_preserve_paragraphs(text)
            # Insert heading boundaries to help chunking
            text = _insert_heading_boundaries(text)
            return text
        except Exception:
            # Fall back to other extractors if PyMuPDF is not available or fails
            pass

    # pypdf extraction
    reader = PdfReader(path)
    texts_pypdf: List[str] = []
    try:
        for page in reader.pages:
            txt = page.extract_text() or ""
            texts_pypdf.append(txt)
    except Exception:
        texts_pypdf = []
    text_pypdf = "\n\n".join(texts_pypdf)
    text_pypdf = _fix_hyphenation(text_pypdf)
    text_pypdf = _normalize_whitespace_preserve_paragraphs(text_pypdf)
    text_pypdf = _insert_heading_boundaries(text_pypdf)

    # Decide if we should try pdfplumber fallback (very sparse output or extremely short)
    needs_fallback = len(text_pypdf.strip()) < 200 or text_pypdf.count("\n") < max(2, len(texts_pypdf) // 4)
    if needs_fallback:
        try:
            import pdfplumber  # type: ignore
            with pdfplumber.open(path) as pdf:
                pages_text = []
                for page in pdf.pages:
                    # Tolerances can help capture columns/tables better
                    t = page.extract_text(x_tolerance=1, y_tolerance=1) or ""
                    pages_text.append(t)
            text_plumb = "\n\n".join(pages_text)
            text_plumb = _fix_hyphenation(text_plumb)
            text_plumb = _normalize_whitespace_preserve_paragraphs(text_plumb)
            text_plumb = _insert_heading_boundaries(text_plumb)
            # Prefer the better (longer, more structured) output
            if len(text_plumb.strip()) > len(text_pypdf.strip()):
                return text_plumb
        except Exception:
            # If pdfplumber fails, keep pypdf output
            pass

    return text_pypdf


def extract_text_from_html(path: str) -> str:
    with open(path, "rb") as f:
        data = f.read()
    soup = BeautifulSoup(data, "html.parser")
    # Remove nav-like elements
    for tag in soup(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_xml(path: str) -> str:
    with open(path, "rb") as f:
        data = f.read()
    soup = BeautifulSoup(data, "xml")
    text = soup.get_text(separator="\n", strip=True)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_csv(path: str) -> str:
    parts: List[str] = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            parts.append(" \t ".join(cell.strip() for cell in row if cell))
    text = "\n".join(parts)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_md(path: str) -> str:
    # Light-weight: treat as plain text (could add markdown parsing if needed)
    return extract_text_from_txt(path)


def extract_text_from_json(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)
        # Convert JSON to a flat text string
        def _flatten(obj) -> List[str]:
            out: List[str] = []
            if isinstance(obj, dict):
                for k, v in obj.items():
                    out.append(str(k))
                    out.extend(_flatten(v))
            elif isinstance(obj, list):
                for it in obj:
                    out.extend(_flatten(it))
            else:
                out.append(str(obj))
            return out
        parts = _flatten(data)
        text = "\n".join(s.strip() for s in parts if s and isinstance(s, str))
        text = _normalize_whitespace_preserve_paragraphs(text)
        return text
    except Exception:
        # Fall back to raw text if not valid JSON
        return extract_text_from_txt(path)


def extract_text_from_docx(path: str) -> str:
    doc = Document(path)
    parts: List[str] = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            parts.append(text)
    text = "\n\n".join(parts)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise ValueError("PPTX support requires optional dependency python-pptx") from e
    prs = Presentation(path)
    parts: List[str] = []
    for slide in prs.slides:
        # Shapes with text
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text:
                    parts.append(shape.text)
            except Exception:
                continue
        # Tables
        for shape in slide.shapes:
            try:
                if not hasattr(shape, 'table') or shape.table is None:
                    continue
                tbl = shape.table
                for row in tbl.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append((cell.text or '').strip())
                    parts.append(" \t ".join(cells))
            except Exception:
                continue
    text = "\n\n".join([p.strip() for p in parts if p and p.strip()])
    return _normalize_whitespace_preserve_paragraphs(text)


def extract_text_from_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as e:
        raise ValueError("XLSX support requires optional dependency openpyxl") from e
    wb = load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = []
            for v in row:
                if v is None:
                    continue
                s = str(v).strip()
                if s:
                    vals.append(s)
            if vals:
                parts.append(" \t ".join(vals))
    text = "\n".join(parts)
    return _normalize_whitespace_preserve_paragraphs(text)


def extract_text_from_image(path: str) -> str:
    try:
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore
    except Exception as e:
        raise ValueError("Image OCR requires optional dependencies pillow and pytesseract") from e
    img = Image.open(path)
    try:
        txt = pytesseract.image_to_string(img)
    except Exception as e:
        raise ValueError(f"OCR failed: {e}") from e
    return _normalize_whitespace_preserve_paragraphs(txt or "")


def extract_text_from_av(path: str, kind: str = "audio") -> str:
    """Transcribe audio/video using Whisper if available; video is first converted to audio via ffmpeg-python.
    kind: 'audio' or 'video'
    """
    try:
        import os as _os
        import tempfile
        import subprocess
        try:
            import whisper  # type: ignore
        except Exception as e:
            raise ValueError("Audio/Video transcription requires optional dependency openai-whisper (whisper)") from e

        src = path
        tmp_audio = None
        if kind == "video":
            # Extract audio to wav using ffmpeg if available
            try:
                import ffmpeg  # type: ignore
                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tf:
                    tmp_audio = tf.name
                (
                    ffmpeg
                    .input(path)
                    .output(tmp_audio, ac=1, ar=16000, format='wav')
                    .overwrite_output()
                    .run(quiet=True)
                )
                src = tmp_audio
            except Exception:
                # Fallback to system ffmpeg if ffmpeg-python is unavailable
                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tf:
                    tmp_audio = tf.name
                try:
                    subprocess.run(["ffmpeg", "-y", "-i", path, "-ac", "1", "-ar", "16000", src], check=True, capture_output=True)
                except Exception as e:
                    raise ValueError("Failed to extract audio from video; install ffmpeg or ffmpeg-python") from e

        # Transcribe with Whisper (uses default/base model; customize as needed)
        try:
            model = whisper.load_model("base")
            result = model.transcribe(src)
            text = (result.get("text") or "").strip()
        finally:
            if tmp_audio and os.path.exists(tmp_audio):
                try:
                    os.remove(tmp_audio)
                except Exception:
                    pass
        return _normalize_whitespace_preserve_paragraphs(text)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Transcription failed: {e}") from e



def _recursive_split(text: str, chunk_size: int, separators: tuple[str, ...]) -> List[str]:
    if not text:
        return []
    if len(text) <= chunk_size or not separators:
        return [text]

    sep = separators[0]
    if sep:
        pieces = text.split(sep)
        rebuilt: List[str] = []
        buf = ""
        joiner = sep
        for piece in pieces:
            candidate = (buf + joiner + piece) if buf else piece
            if len(candidate) <= chunk_size:
                buf = candidate
            else:
                if buf:
                    rebuilt.append(buf)
                if len(piece) <= chunk_size:
                    buf = piece
                else:
                    rebuilt.extend(_recursive_split(piece, chunk_size, separators[1:]))
                    buf = ""
        if buf:
            rebuilt.append(buf)
        return rebuilt
    else:
        return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]


def _apply_overlap(chunks: List[str], overlap: int) -> List[str]:
    if overlap <= 0 or not chunks:
        return chunks
    out: List[str] = []
    prev_tail = ""
    for ch in chunks:
        prefix = prev_tail
        combined = (prefix + ch) if prefix else ch
        out.append(combined)
        prev_tail = ch[-overlap:]
    return out


def chunk_text(text: str, params: ChunkParams = ChunkParams()) -> List[str]:
    """Split text into chunks with optional adaptive sizing.
    When CHUNK_AUTO_TUNE=true, choose a chunk size based on paragraph statistics
    clamped to [CHUNK_MIN_SIZE, CHUNK_MAX_SIZE], and set overlap proportionally
    using CHUNK_OVERLAP_RATIO.
    """
    # Normalize while preserving paragraph boundaries; add extra spacing around likely headings
    text = _normalize_whitespace_preserve_paragraphs(text)
    text = _insert_heading_boundaries(text)

    # Start with configured defaults
    eff_chunk_size = int(params.chunk_size)
    eff_overlap = int(params.chunk_overlap)

    # Adaptive tuning based on paragraph distribution
    try:
        if getattr(settings, "chunk_auto_tune", False):
            paras = [p for p in re.split(PARA_SPLIT_RE, text) if p and p.strip()]
            if paras:
                import statistics
                lens = [len(p) for p in paras]
                med = int(statistics.median(lens)) if lens else eff_chunk_size
                avg = int(sum(lens) / max(1, len(lens))) if lens else eff_chunk_size
                # Heuristic target: 2.5x median or 2x average, whichever larger, but within min/max
                target = max(
                    int(getattr(settings, "chunk_min_size", 800)),
                    min(
                        int(getattr(settings, "chunk_max_size", 3500)),
                        int(max(med * 2.5, avg * 2.0, eff_chunk_size)),
                    ),
                )
                eff_chunk_size = int(target)
                # Overlap as a ratio of chosen size (bounded)
                ratio = float(getattr(settings, "chunk_overlap_ratio", 0.1) or 0.0)
                eff_overlap = int(min(max(0, ratio * eff_chunk_size), max(600, eff_chunk_size // 3)))
    except Exception:
        # Fall back to configured params on any issues
        eff_chunk_size = int(params.chunk_size)
        eff_overlap = int(params.chunk_overlap)

    base_chunks = _recursive_split(text, eff_chunk_size, params.separators)
    if not base_chunks:
        return []
    return _apply_overlap(base_chunks, eff_overlap)